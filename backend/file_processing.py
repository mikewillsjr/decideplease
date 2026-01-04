"""File processing utilities for DecidePlease file uploads."""

import base64
import io
from typing import Dict, List, Any, Optional, Tuple
from PIL import Image

from .openrouter import query_model
from .config import DESCRIPTION_MODEL

# File type configurations
ALLOWED_MIME_TYPES = {
    # Images
    'image/jpeg': 'image',
    'image/png': 'image',
    'image/gif': 'image',
    'image/webp': 'image',
    # PDF
    'application/pdf': 'pdf',
    # Office documents
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'docx',
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': 'xlsx',
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'pptx',
}

# Magic bytes for file type validation
MAGIC_BYTES = {
    'image/jpeg': [b'\xff\xd8\xff'],
    'image/png': [b'\x89PNG\r\n\x1a\n'],
    'image/gif': [b'GIF87a', b'GIF89a'],
    'image/webp': [b'RIFF'],  # WebP starts with RIFF....WEBP
    'application/pdf': [b'%PDF'],
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': [b'PK\x03\x04'],
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': [b'PK\x03\x04'],
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': [b'PK\x03\x04'],
}

MAX_FILE_SIZE = 10 * 1024 * 1024  # 10MB
MAX_FILES = 5

# Prompt for Gemini Flash to describe images for text-only models
IMAGE_DESCRIPTION_PROMPT = """Describe this image in detail for someone who cannot see it.
Include:
- What type of image it is (photo, diagram, chart, screenshot, document, etc.)
- Key visual elements and their arrangement
- Any text visible in the image (transcribe it if readable)
- Colors, shapes, and composition if relevant
- Overall context and meaning

Be thorough but concise. This description will be used by an AI that cannot see images to help answer a user's question about this image."""


class FileValidationError(Exception):
    """Raised when file validation fails."""
    pass


def validate_file(filename: str, content_type: str, data: bytes) -> str:
    """
    Validate a file's type and size.

    Args:
        filename: Original filename
        content_type: MIME type
        data: Raw file bytes

    Returns:
        File category ('image', 'pdf', 'docx', 'xlsx', 'pptx')

    Raises:
        FileValidationError: If validation fails
    """
    # Check MIME type
    if content_type not in ALLOWED_MIME_TYPES:
        raise FileValidationError(
            f"File type '{content_type}' is not supported. "
            "Allowed: JPEG, PNG, GIF, WebP, PDF, DOCX, XLSX, PPTX"
        )

    # Check size
    if len(data) > MAX_FILE_SIZE:
        size_mb = len(data) / (1024 * 1024)
        raise FileValidationError(
            f"File '{filename}' is {size_mb:.1f}MB, exceeds 10MB limit"
        )

    # Validate magic bytes
    if content_type in MAGIC_BYTES:
        valid = False
        for magic in MAGIC_BYTES[content_type]:
            if data.startswith(magic):
                valid = True
                break

        # Special handling for WebP (RIFF....WEBP format)
        if content_type == 'image/webp' and data.startswith(b'RIFF'):
            if len(data) >= 12 and data[8:12] == b'WEBP':
                valid = True
            else:
                valid = False

        if not valid:
            raise FileValidationError(
                f"File '{filename}' content doesn't match declared type '{content_type}'"
            )

    return ALLOWED_MIME_TYPES[content_type]


def validate_files(files: List[Dict[str, Any]]) -> None:
    """
    Validate a list of file attachments.

    Args:
        files: List of file dicts with 'filename', 'content_type', 'data' (base64)

    Raises:
        FileValidationError: If validation fails
    """
    if len(files) > MAX_FILES:
        raise FileValidationError(f"Maximum {MAX_FILES} files allowed, got {len(files)}")

    for file in files:
        try:
            data = base64.b64decode(file['data'])
        except Exception:
            raise FileValidationError(f"Invalid base64 data for file '{file['filename']}'")

        validate_file(file['filename'], file['content_type'], data)


def process_image_to_data_uri(data: bytes, content_type: str) -> str:
    """
    Convert image bytes to a data URI.

    Args:
        data: Raw image bytes
        content_type: MIME type

    Returns:
        Data URI string (data:image/jpeg;base64,...)
    """
    b64 = base64.b64encode(data).decode('utf-8')
    return f"data:{content_type};base64,{b64}"


def extract_pdf_text(data: bytes) -> str:
    """
    Extract text content from a PDF file.

    Args:
        data: Raw PDF bytes

    Returns:
        Extracted text content
    """
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    text_parts = []

    for i, page in enumerate(reader.pages, start=1):
        page_text = page.extract_text()
        if page_text.strip():
            text_parts.append(f"[Page {i}]\n{page_text}")

    if not text_parts:
        return "[PDF contains no extractable text - may be image-based]"

    return "\n\n".join(text_parts)


def extract_docx_text(data: bytes) -> str:
    """
    Extract text content from a Word document.

    Args:
        data: Raw DOCX bytes

    Returns:
        Extracted text content
    """
    from docx import Document

    doc = Document(io.BytesIO(data))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

    if not paragraphs:
        return "[Document contains no extractable text]"

    return "\n\n".join(paragraphs)


def extract_xlsx_text(data: bytes) -> str:
    """
    Extract text content from an Excel spreadsheet.

    Args:
        data: Raw XLSX bytes

    Returns:
        Extracted text content (formatted as tables)
    """
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), data_only=True)
    text_parts = []

    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        rows = []

        for row in sheet.iter_rows(values_only=True):
            # Filter out completely empty rows
            if any(cell is not None for cell in row):
                row_str = " | ".join(str(cell) if cell is not None else "" for cell in row)
                rows.append(row_str)

        if rows:
            text_parts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))

    if not text_parts:
        return "[Spreadsheet contains no data]"

    return "\n\n".join(text_parts)


def extract_pptx_text(data: bytes) -> str:
    """
    Extract text content from a PowerPoint presentation.

    Args:
        data: Raw PPTX bytes

    Returns:
        Extracted text content
    """
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    text_parts = []

    for i, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text)

        if slide_texts:
            text_parts.append(f"[Slide {i}]\n" + "\n".join(slide_texts))

    if not text_parts:
        return "[Presentation contains no extractable text]"

    return "\n\n".join(text_parts)


async def generate_image_description(data_uri: str, filename: str) -> str:
    """
    Use Gemini Flash to generate a text description of an image.

    This is used for text-only models (like DeepSeek) that can't process images.

    Args:
        data_uri: Image data URI (data:image/jpeg;base64,...)
        filename: Original filename for context

    Returns:
        Text description of the image
    """
    messages = [{
        "role": "user",
        "content": [
            {"type": "text", "text": IMAGE_DESCRIPTION_PROMPT},
            {"type": "image_url", "image_url": {"url": data_uri}}
        ]
    }]

    response = await query_model(DESCRIPTION_MODEL, messages, timeout=30.0)

    if response is None:
        return f"[Image: {filename} - description unavailable]"

    description = response.get('content', '')
    if not description:
        return f"[Image: {filename} - description unavailable]"

    return description


async def process_files(files: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    """
    Process uploaded files for council consumption.

    Args:
        files: List of file dicts with 'filename', 'content_type', 'data' (base64)

    Returns:
        List of processed file dicts with:
        - filename: Original filename
        - file_type: Category ('image', 'pdf', 'docx', 'xlsx', 'pptx')
        - content_type: MIME type
        - data_uri: For images, the data URI
        - extracted_text: For documents, the extracted text
    """
    processed = []

    for file in files:
        data = base64.b64decode(file['data'])
        file_type = ALLOWED_MIME_TYPES[file['content_type']]

        result = {
            'filename': file['filename'],
            'file_type': file_type,
            'content_type': file['content_type'],
        }

        if file_type == 'image':
            result['data_uri'] = process_image_to_data_uri(data, file['content_type'])
        elif file_type == 'pdf':
            result['extracted_text'] = extract_pdf_text(data)
        elif file_type == 'docx':
            result['extracted_text'] = extract_docx_text(data)
        elif file_type == 'xlsx':
            result['extracted_text'] = extract_xlsx_text(data)
        elif file_type == 'pptx':
            result['extracted_text'] = extract_pptx_text(data)

        processed.append(result)

    return processed


async def generate_image_descriptions_for_text_models(
    processed_files: List[Dict[str, Any]]
) -> Dict[str, str]:
    """
    Generate text descriptions for all images using Gemini Flash.

    This is used so text-only models (like DeepSeek) can understand images.

    Args:
        processed_files: List of processed file dicts

    Returns:
        Dict mapping filename to description
    """
    import asyncio

    descriptions = {}
    image_files = [f for f in processed_files if f['file_type'] == 'image']

    if not image_files:
        return descriptions

    # Generate descriptions in parallel
    tasks = [
        generate_image_description(f['data_uri'], f['filename'])
        for f in image_files
    ]

    results = await asyncio.gather(*tasks)

    for file, description in zip(image_files, results):
        descriptions[file['filename']] = description

    return descriptions
